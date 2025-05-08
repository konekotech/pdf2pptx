import os
import tempfile
from cairosvg import svg2png
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt
import sys

def pdf_to_pptx(pdf_path, output_pptx):
    pages = convert_from_path(pdf_path)
    prs = Presentation()

    # 最初のページのサイズでアスペクト比取得
    first_page = pages[0]
    pdf_width_px, pdf_height_px = first_page.size
    aspect_ratio = pdf_width_px / pdf_height_px

    # 任意の高さを設定（ポイント単位、ここでは7.5インチ = 540ポイントに）
    slide_height = Inches(7.5)
    slide_width = Pt(aspect_ratio * slide_height.pt)

    prs.slide_width = slide_width
    prs.slide_height = slide_height

    for i, page in enumerate(pages):
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
            page.save(tmp_img.name, "PNG")
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide

            slide.shapes.add_picture(tmp_img.name, 0, 0, width=slide_width, height=slide_height)
            img_path = tmp_img.name

        os.remove(img_path)

    prs.save(output_pptx)
    print(f"PowerPoint saved to {output_pptx}")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python main.py <input_pdf> <output_pptx>")
        sys.exit(1)

    input_pdf = sys.argv[1]
    output_pptx = sys.argv[2]

    if not os.path.exists(input_pdf):
        print(f"Error: {input_pdf} does not exist.")
        sys.exit(1)

    pdf_to_pptx(input_pdf, output_pptx)


